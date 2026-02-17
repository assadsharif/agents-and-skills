"""
PPTX MCP Server — PowerPoint presentation creation, editing, and analysis toolkit.

TOOLS:
    pptx_generate_extraction_code    Generate text/content extraction code
    pptx_generate_creation_code      Generate presentation creation code
    pptx_generate_slide_html         Generate HTML for slide content
    pptx_generate_template_code      Generate template-based workflow code
    pptx_generate_edit_code          Generate OOXML editing code
    pptx_generate_color_palette      Generate color palette recommendations
    pptx_generate_layout             Generate layout configurations
    pptx_generate_thumbnail_code     Generate thumbnail/validation code
    pptx_detect_antipatterns         Detect anti-patterns in PPTX code
    pptx_generate_scaffold           Generate complete PPTX project scaffold
"""

import json
import re
from typing import Optional

from mcp.server.fastmcp import FastMCP
from pydantic import BaseModel, ConfigDict, Field, field_validator

# ---------------------------------------------------------------------------
# FastMCP server instance
# ---------------------------------------------------------------------------

mcp = FastMCP("pptx_mcp")

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

EXTRACTION_TYPES = ["text", "xml", "metadata", "images"]
CREATION_TYPES = ["basic", "html2pptx", "pptxgenjs"]
LAYOUT_TYPES = ["title", "content", "two_column", "image_text", "quote", "section"]
ASPECT_RATIOS = ["16:9", "4:3", "16:10"]

# Web-safe fonts for PowerPoint
WEB_SAFE_FONTS = [
    "Arial",
    "Helvetica",
    "Times New Roman",
    "Georgia",
    "Courier New",
    "Verdana",
    "Tahoma",
    "Trebuchet MS",
    "Impact",
]

# Color palettes from skill documentation
COLOR_PALETTES = {
    "classic_blue": {
        "name": "Classic Blue",
        "colors": {
            "primary": "#1C2833",
            "secondary": "#2E4053",
            "accent": "#AAB7B8",
            "background": "#F4F6F6",
        },
        "mood": "professional, corporate, trustworthy",
    },
    "teal_coral": {
        "name": "Teal & Coral",
        "colors": {
            "primary": "#5EA8A7",
            "secondary": "#277884",
            "accent": "#FE4447",
            "background": "#FFFFFF",
        },
        "mood": "modern, vibrant, creative",
    },
    "bold_red": {
        "name": "Bold Red",
        "colors": {
            "primary": "#C0392B",
            "secondary": "#E74C3C",
            "accent": "#F39C12",
            "background": "#FFFFFF",
        },
        "mood": "energetic, urgent, passionate",
    },
    "warm_blush": {
        "name": "Warm Blush",
        "colors": {
            "primary": "#A49393",
            "secondary": "#EED6D3",
            "accent": "#E8B4B8",
            "background": "#FAF7F2",
        },
        "mood": "elegant, feminine, soft",
    },
    "burgundy_luxury": {
        "name": "Burgundy Luxury",
        "colors": {
            "primary": "#5D1D2E",
            "secondary": "#951233",
            "accent": "#997929",
            "background": "#FFFFFF",
        },
        "mood": "luxurious, sophisticated, premium",
    },
    "deep_purple": {
        "name": "Deep Purple & Emerald",
        "colors": {
            "primary": "#B165FB",
            "secondary": "#181B24",
            "accent": "#40695B",
            "background": "#FFFFFF",
        },
        "mood": "creative, bold, modern",
    },
    "forest_cream": {
        "name": "Cream & Forest Green",
        "colors": {
            "primary": "#40695B",
            "secondary": "#FFE1C7",
            "accent": "#40695B",
            "background": "#FCFCFC",
        },
        "mood": "natural, calm, organic",
    },
    "black_gold": {
        "name": "Black & Gold",
        "colors": {
            "primary": "#000000",
            "secondary": "#BF9A4A",
            "accent": "#BF9A4A",
            "background": "#F4F6F6",
        },
        "mood": "premium, elegant, powerful",
    },
    "sage_terracotta": {
        "name": "Sage & Terracotta",
        "colors": {
            "primary": "#87A96B",
            "secondary": "#E07A5F",
            "accent": "#2C2C2C",
            "background": "#F4F1DE",
        },
        "mood": "earthy, warm, natural",
    },
    "charcoal_red": {
        "name": "Charcoal & Red",
        "colors": {
            "primary": "#292929",
            "secondary": "#E33737",
            "accent": "#CCCBCB",
            "background": "#FFFFFF",
        },
        "mood": "bold, modern, high-contrast",
    },
}

# Code templates
TEXT_EXTRACTION_TEMPLATE = '''"""Extract text from PowerPoint using markitdown."""
import subprocess
from pathlib import Path

def extract_text(pptx_path: str, output_path: str | None = None) -> str:
    """Extract text content from PowerPoint presentation.

    Args:
        pptx_path: Path to input PPTX file
        output_path: Optional path to save extracted text

    Returns:
        Extracted text content as markdown
    """
    result = subprocess.run(
        ["python", "-m", "markitdown", pptx_path],
        capture_output=True,
        text=True,
        check=True,
    )

    text = result.stdout

    if output_path:
        Path(output_path).write_text(text, encoding="utf-8")

    return text
'''

XML_EXTRACTION_TEMPLATE = '''"""Extract raw XML from PowerPoint for detailed analysis."""
import subprocess
from pathlib import Path
import xml.etree.ElementTree as ET

def unpack_pptx(pptx_path: str, output_dir: str) -> None:
    """Unpack PPTX to access raw XML files.

    Args:
        pptx_path: Path to input PPTX file
        output_dir: Directory to extract contents
    """
    subprocess.run(
        ["python", "ooxml/scripts/unpack.py", pptx_path, output_dir],
        check=True,
    )

def read_slide_xml(unpacked_dir: str, slide_num: int) -> str:
    """Read raw XML content of a specific slide.

    Args:
        unpacked_dir: Path to unpacked PPTX directory
        slide_num: Slide number (1-based)

    Returns:
        XML content as string
    """
    slide_path = Path(unpacked_dir) / f"ppt/slides/slide{slide_num}.xml"
    return slide_path.read_text(encoding="utf-8")

def extract_theme_colors(unpacked_dir: str) -> dict:
    """Extract color scheme from theme file.

    Args:
        unpacked_dir: Path to unpacked PPTX directory

    Returns:
        Dictionary of theme colors
    """
    theme_path = Path(unpacked_dir) / "ppt/theme/theme1.xml"
    tree = ET.parse(theme_path)
    root = tree.getroot()

    # Extract color scheme (simplified)
    colors = {}
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    for clr in root.findall(".//a:clrScheme/*", ns):
        tag = clr.tag.split("}")[-1]
        srgb = clr.find(".//a:srgbClr", ns)
        if srgb is not None:
            colors[tag] = srgb.get("val")

    return colors
'''

HTML2PPTX_TEMPLATE = '''"""Create PowerPoint from HTML slides using html2pptx."""
const {{ html2pptx }} = require('./html2pptx.js');
const PptxGenJS = require('pptxgenjs');
const fs = require('fs');

async function createPresentation(htmlFiles, outputPath) {{
    const pptx = new PptxGenJS();

    // Set presentation properties
    pptx.layout = '{aspect_ratio}';
    pptx.title = '{title}';
    pptx.author = '{author}';

    // Process each HTML slide
    for (const htmlFile of htmlFiles) {{
        const htmlContent = fs.readFileSync(htmlFile, 'utf8');
        const slide = pptx.addSlide();

        // Convert HTML to PPTX elements
        const elements = await html2pptx(htmlContent);

        for (const element of elements) {{
            if (element.type === 'text') {{
                slide.addText(element.text, element.options);
            }} else if (element.type === 'shape') {{
                slide.addShape(element.shape, element.options);
            }} else if (element.type === 'image') {{
                slide.addImage(element.options);
            }}
        }}
    }}

    await pptx.writeFile(outputPath);
    console.log(`Presentation saved to ${{outputPath}}`);
}}

// Usage
const htmlFiles = ['slide1.html', 'slide2.html', 'slide3.html'];
createPresentation(htmlFiles, 'output.pptx');
'''

SLIDE_HTML_TEMPLATE = """<!DOCTYPE html>
<html>
<head>
<style>
body {{
    margin: 0;
    padding: 0;
    width: {width}pt;
    height: {height}pt;
    font-family: {font_family};
    background: {background};
    display: flex;
    flex-direction: column;
}}
{additional_styles}
</style>
</head>
<body>
{content}
</body>
</html>
"""

TEMPLATE_WORKFLOW_CODE = '''"""Template-based PowerPoint workflow."""
import subprocess
import json
from pathlib import Path

def analyze_template(template_path: str, output_dir: str) -> dict:
    """Analyze template and create inventory.

    Args:
        template_path: Path to template PPTX
        output_dir: Working directory

    Returns:
        Template analysis results
    """
    # Extract text
    subprocess.run(
        ["python", "-m", "markitdown", template_path],
        capture_output=True,
        text=True,
    )

    # Create thumbnails
    subprocess.run(
        ["python", "scripts/thumbnail.py", template_path, f"{output_dir}/thumbnails"],
        check=True,
    )

    return {"template": template_path, "thumbnails": f"{output_dir}/thumbnails.jpg"}

def rearrange_slides(template_path: str, output_path: str, slide_order: list[int]) -> None:
    """Rearrange slides from template.

    Args:
        template_path: Path to template PPTX
        output_path: Path for output PPTX
        slide_order: List of slide indices (0-based)
    """
    order_str = ",".join(str(i) for i in slide_order)
    subprocess.run(
        ["python", "scripts/rearrange.py", template_path, output_path, order_str],
        check=True,
    )

def extract_inventory(pptx_path: str, output_json: str) -> dict:
    """Extract text inventory from presentation.

    Args:
        pptx_path: Path to PPTX file
        output_json: Path for output JSON

    Returns:
        Inventory data
    """
    subprocess.run(
        ["python", "scripts/inventory.py", pptx_path, output_json],
        check=True,
    )

    return json.loads(Path(output_json).read_text())

def replace_text(pptx_path: str, replacements_json: str, output_path: str) -> None:
    """Replace text in presentation.

    Args:
        pptx_path: Path to input PPTX
        replacements_json: Path to replacements JSON
        output_path: Path for output PPTX
    """
    subprocess.run(
        ["python", "scripts/replace.py", pptx_path, replacements_json, output_path],
        check=True,
    )
'''

OOXML_EDIT_TEMPLATE = '''"""Edit PowerPoint via OOXML manipulation."""
import subprocess
from pathlib import Path
import xml.etree.ElementTree as ET

# Register namespaces
NAMESPACES = {{
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}}

for prefix, uri in NAMESPACES.items():
    ET.register_namespace(prefix, uri)

def unpack_pptx(pptx_path: str, output_dir: str) -> None:
    """Unpack PPTX for editing."""
    subprocess.run(
        ["python", "ooxml/scripts/unpack.py", pptx_path, output_dir],
        check=True,
    )

def pack_pptx(input_dir: str, output_path: str) -> None:
    """Pack directory back to PPTX."""
    subprocess.run(
        ["python", "ooxml/scripts/pack.py", input_dir, output_path],
        check=True,
    )

def validate_pptx(unpacked_dir: str, original_path: str) -> bool:
    """Validate unpacked PPTX structure."""
    result = subprocess.run(
        ["python", "ooxml/scripts/validate.py", unpacked_dir, "--original", original_path],
        capture_output=True,
        text=True,
    )
    return result.returncode == 0

def edit_slide_text(unpacked_dir: str, slide_num: int, shape_id: str, new_text: str) -> None:
    """Edit text in a specific shape.

    Args:
        unpacked_dir: Path to unpacked PPTX
        slide_num: Slide number (1-based)
        shape_id: Shape identifier
        new_text: New text content
    """
    slide_path = Path(unpacked_dir) / f"ppt/slides/slide{{slide_num}}.xml"
    tree = ET.parse(slide_path)
    root = tree.getroot()

    # Find and update text (simplified)
    for t_elem in root.findall(".//a:t", NAMESPACES):
        if t_elem.text:
            # Update logic based on shape_id
            pass

    tree.write(slide_path, xml_declaration=True, encoding="UTF-8")
'''

THUMBNAIL_CODE = '''"""Generate thumbnails for visual validation."""
import subprocess
from pathlib import Path

def create_thumbnails(pptx_path: str, output_prefix: str, cols: int = 5) -> str:
    """Create thumbnail grid of presentation.

    Args:
        pptx_path: Path to PPTX file
        output_prefix: Prefix for output files
        cols: Number of columns (3-6)

    Returns:
        Path to generated thumbnail image
    """
    subprocess.run(
        ["python", "scripts/thumbnail.py", pptx_path, output_prefix, "--cols", str(cols)],
        check=True,
    )

    return f"{output_prefix}.jpg"

def convert_to_pdf(pptx_path: str) -> str:
    """Convert PPTX to PDF for image extraction.

    Args:
        pptx_path: Path to PPTX file

    Returns:
        Path to generated PDF
    """
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf", pptx_path],
        check=True,
    )

    return pptx_path.replace(".pptx", ".pdf")

def pdf_to_images(pdf_path: str, output_prefix: str, dpi: int = 150) -> list[str]:
    """Convert PDF pages to images.

    Args:
        pdf_path: Path to PDF file
        output_prefix: Prefix for output images
        dpi: Resolution (default 150)

    Returns:
        List of generated image paths
    """
    subprocess.run(
        ["pdftoppm", "-jpeg", "-r", str(dpi), pdf_path, output_prefix],
        check=True,
    )

    # Find generated files
    parent = Path(pdf_path).parent
    return sorted(str(p) for p in parent.glob(f"{output_prefix}-*.jpg"))
'''

# Anti-patterns
ANTIPATTERNS = [
    {
        "pattern": r"<div>[^<]*text[^<]*</div>",
        "issue": "Text directly in div",
        "fix": "Wrap text in <p>, <h1>-<h6>, <ul>, or <ol> tags",
    },
    {
        "pattern": r"linear-gradient|radial-gradient",
        "issue": "CSS gradients won't convert",
        "fix": "Rasterize gradients to PNG with Sharp first",
    },
    {
        "pattern": r"font-family:\s*['\"]?(?!Arial|Helvetica|Times|Georgia|Courier|Verdana|Tahoma|Trebuchet|Impact)",
        "issue": "Non-web-safe font",
        "fix": "Use web-safe fonts: Arial, Helvetica, Times New Roman, Georgia, Courier New, Verdana, Tahoma, Trebuchet MS, Impact",
    },
    {
        "pattern": r"[•\-\*]\s+\w",
        "issue": "Manual bullet symbols",
        "fix": "Use <ul> or <ol> lists instead of manual bullet characters",
    },
    {
        "pattern": r"<span>[^<]+</span>(?!</p>)",
        "issue": "Text in span without parent",
        "fix": "Wrap in <p> tag: <p><span>text</span></p>",
    },
    {
        "pattern": r"slide\d+\.xml.*slide\d+\.xml",
        "issue": "Multiple slide edits without validation",
        "fix": "Validate after each edit with validate.py",
    },
    {
        "pattern": r"writeFile.*without.*validation",
        "issue": "Saving without visual validation",
        "fix": "Generate thumbnails and inspect before finalizing",
    },
]

# ---------------------------------------------------------------------------
# Pydantic models
# ---------------------------------------------------------------------------

_CFG = ConfigDict(str_strip_whitespace=True, extra="forbid")


class ExtractionInput(BaseModel):
    model_config = _CFG
    extraction_type: str = Field(..., description="Type: text, xml, metadata, images")

    @field_validator("extraction_type")
    @classmethod
    def _check_type(cls, v: str) -> str:
        if v not in EXTRACTION_TYPES:
            raise ValueError(f"extraction_type must be one of {EXTRACTION_TYPES}")
        return v


class CreationInput(BaseModel):
    model_config = _CFG
    creation_type: str = Field(..., description="Type: basic, html2pptx, pptxgenjs")
    aspect_ratio: str = Field(
        default="16:9", description="Aspect ratio: 16:9, 4:3, 16:10"
    )
    title: str = Field(default="Presentation", description="Presentation title")
    author: str = Field(default="", description="Author name")

    @field_validator("creation_type")
    @classmethod
    def _check_type(cls, v: str) -> str:
        if v not in CREATION_TYPES:
            raise ValueError(f"creation_type must be one of {CREATION_TYPES}")
        return v

    @field_validator("aspect_ratio")
    @classmethod
    def _check_ratio(cls, v: str) -> str:
        if v not in ASPECT_RATIOS:
            raise ValueError(f"aspect_ratio must be one of {ASPECT_RATIOS}")
        return v


class SlideHtmlInput(BaseModel):
    model_config = _CFG
    layout_type: str = Field(
        ..., description="Type: title, content, two_column, image_text, quote, section"
    )
    aspect_ratio: str = Field(default="16:9", description="Aspect ratio")
    title: str = Field(default="", description="Slide title")
    content: str = Field(default="", description="Main content")
    palette_name: str = Field(default="classic_blue", description="Color palette name")

    @field_validator("layout_type")
    @classmethod
    def _check_type(cls, v: str) -> str:
        if v not in LAYOUT_TYPES:
            raise ValueError(f"layout_type must be one of {LAYOUT_TYPES}")
        return v


class TemplateInput(BaseModel):
    model_config = _CFG
    workflow_step: str = Field(
        ..., description="Step: analyze, rearrange, inventory, replace"
    )

    @field_validator("workflow_step")
    @classmethod
    def _check_step(cls, v: str) -> str:
        if v not in ("analyze", "rearrange", "inventory", "replace", "full"):
            raise ValueError(
                "workflow_step must be analyze, rearrange, inventory, replace, or full"
            )
        return v


class EditInput(BaseModel):
    model_config = _CFG
    edit_type: str = Field(..., description="Type: text, shape, image, notes")

    @field_validator("edit_type")
    @classmethod
    def _check_type(cls, v: str) -> str:
        if v not in ("text", "shape", "image", "notes"):
            raise ValueError("edit_type must be text, shape, image, or notes")
        return v


class ColorPaletteInput(BaseModel):
    model_config = _CFG
    mood: str = Field(default="", description="Desired mood/style")
    industry: str = Field(default="", description="Industry/domain")


class LayoutInput(BaseModel):
    model_config = _CFG
    layout_type: str = Field(..., description="Layout type")
    aspect_ratio: str = Field(default="16:9", description="Aspect ratio")

    @field_validator("layout_type")
    @classmethod
    def _check_type(cls, v: str) -> str:
        if v not in LAYOUT_TYPES:
            raise ValueError(f"layout_type must be one of {LAYOUT_TYPES}")
        return v


class ThumbnailInput(BaseModel):
    model_config = _CFG
    include_pdf_conversion: bool = Field(
        default=False, description="Include PDF conversion code"
    )


class AntipatternInput(BaseModel):
    model_config = _CFG
    code: str = Field(
        ..., min_length=10, max_length=50000, description="Code to analyze"
    )


class ScaffoldInput(BaseModel):
    model_config = _CFG
    project_name: str = Field(
        ..., min_length=1, max_length=100, description="Project name"
    )
    workflow: str = Field(
        default="html2pptx", description="Workflow: html2pptx, template, ooxml"
    )

    @field_validator("project_name")
    @classmethod
    def _check_name(cls, v: str) -> str:
        if not re.match(r"^[a-zA-Z][a-zA-Z0-9_-]*$", v):
            raise ValueError(
                "project_name must start with letter, contain only alphanumeric, underscore, hyphen"
            )
        return v


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------


def _get_dimensions(aspect_ratio: str) -> tuple[int, int]:
    """Get slide dimensions for aspect ratio."""
    dims = {
        "16:9": (720, 405),
        "4:3": (720, 540),
        "16:10": (720, 450),
    }
    return dims.get(aspect_ratio, (720, 405))


def _detect_antipatterns(code: str) -> list[dict]:
    """Detect anti-patterns in PPTX code."""
    findings = []
    for ap in ANTIPATTERNS:
        if re.search(ap["pattern"], code, re.IGNORECASE):
            findings.append(
                {
                    "issue": ap["issue"],
                    "fix": ap["fix"],
                    "pattern": ap["pattern"],
                }
            )
    return findings


def _generate_slide_html(
    layout_type: str, aspect_ratio: str, title: str, content: str, palette: dict
) -> str:
    """Generate HTML for a slide."""
    width, height = _get_dimensions(aspect_ratio)
    colors = palette["colors"]

    if layout_type == "title":
        html_content = f"""<div style="flex: 1; display: flex; flex-direction: column; justify-content: center; align-items: center; padding: 40pt;">
    <h1 style="font-size: 48pt; color: {colors['primary']}; margin: 0; text-align: center;">{title}</h1>
    <p style="font-size: 24pt; color: {colors['secondary']}; margin-top: 20pt; text-align: center;">{content}</p>
</div>"""
    elif layout_type == "content":
        html_content = f"""<div style="padding: 30pt;">
    <h2 style="font-size: 36pt; color: {colors['primary']}; margin: 0 0 20pt 0;">{title}</h2>
    <p style="font-size: 18pt; color: {colors['secondary']}; line-height: 1.5;">{content}</p>
</div>"""
    elif layout_type == "two_column":
        html_content = f"""<div style="padding: 30pt;">
    <h2 style="font-size: 32pt; color: {colors['primary']}; margin: 0 0 20pt 0;">{title}</h2>
    <div style="display: flex; gap: 30pt;">
        <div style="flex: 1;">
            <p style="font-size: 16pt; color: {colors['secondary']};">{content}</p>
        </div>
        <div style="flex: 1;" class="placeholder">
            <!-- Chart/image placeholder -->
        </div>
    </div>
</div>"""
    elif layout_type == "quote":
        html_content = f"""<div style="flex: 1; display: flex; flex-direction: column; justify-content: center; padding: 60pt; background: {colors['primary']};">
    <p style="font-size: 28pt; color: {colors['background']}; font-style: italic; text-align: center;">"{content}"</p>
    <p style="font-size: 18pt; color: {colors['accent']}; text-align: center; margin-top: 20pt;">— {title}</p>
</div>"""
    elif layout_type == "section":
        html_content = f"""<div style="flex: 1; display: flex; align-items: center; justify-content: center; background: {colors['primary']};">
    <h1 style="font-size: 56pt; color: {colors['background']}; text-align: center;">{title}</h1>
</div>"""
    else:  # image_text
        html_content = f"""<div style="display: flex; height: 100%;">
    <div style="flex: 1; padding: 30pt; display: flex; flex-direction: column; justify-content: center;">
        <h2 style="font-size: 32pt; color: {colors['primary']}; margin: 0 0 20pt 0;">{title}</h2>
        <p style="font-size: 16pt; color: {colors['secondary']}; line-height: 1.5;">{content}</p>
    </div>
    <div style="flex: 1;" class="placeholder">
        <!-- Image placeholder -->
    </div>
</div>"""

    return SLIDE_HTML_TEMPLATE.format(
        width=width,
        height=height,
        font_family="Arial, sans-serif",
        background=colors["background"],
        additional_styles="",
        content=html_content,
    )


# ---------------------------------------------------------------------------
# MCP tools
# ---------------------------------------------------------------------------


@mcp.tool()
async def pptx_generate_extraction_code(
    extraction_type: str,
) -> str:
    """Generate code for extracting content from PowerPoint. Returns Python code string."""
    try:
        inp = ExtractionInput(extraction_type=extraction_type)

        if inp.extraction_type == "text":
            code = TEXT_EXTRACTION_TEMPLATE
        elif inp.extraction_type == "xml":
            code = XML_EXTRACTION_TEMPLATE
        elif inp.extraction_type == "metadata":
            code = '''"""Extract metadata from PowerPoint."""
from pypptx import Presentation

def extract_metadata(pptx_path: str) -> dict:
    """Extract presentation metadata.

    Args:
        pptx_path: Path to PPTX file

    Returns:
        Dictionary of metadata
    """
    prs = Presentation(pptx_path)
    core = prs.core_properties

    return {
        "title": core.title,
        "author": core.author,
        "subject": core.subject,
        "keywords": core.keywords,
        "created": str(core.created) if core.created else None,
        "modified": str(core.modified) if core.modified else None,
        "slide_count": len(prs.slides),
    }
'''
        else:  # images
            code = '''"""Extract images from PowerPoint slides."""
from pathlib import Path
import subprocess

def extract_images(pptx_path: str, output_dir: str) -> list[str]:
    """Extract all images from presentation.

    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory for extracted images

    Returns:
        List of extracted image paths
    """
    # Unpack PPTX
    subprocess.run(
        ["python", "ooxml/scripts/unpack.py", pptx_path, output_dir],
        check=True,
    )

    # Find media files
    media_dir = Path(output_dir) / "ppt" / "media"
    if not media_dir.exists():
        return []

    return [str(p) for p in media_dir.iterdir() if p.is_file()]
'''

        return json.dumps({"code": code, "type": inp.extraction_type})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def pptx_generate_creation_code(
    creation_type: str,
    aspect_ratio: str = "16:9",
    title: str = "Presentation",
    author: str = "",
) -> str:
    """Generate code for creating PowerPoint presentations. Returns code string."""
    try:
        inp = CreationInput(
            creation_type=creation_type,
            aspect_ratio=aspect_ratio,
            title=title,
            author=author,
        )

        if inp.creation_type == "html2pptx":
            code = HTML2PPTX_TEMPLATE.format(
                aspect_ratio=inp.aspect_ratio.replace(":", "_"),
                title=inp.title,
                author=inp.author,
            )
        elif inp.creation_type == "pptxgenjs":
            code = f"""/**
 * Create PowerPoint using PptxGenJS directly.
 */
const PptxGenJS = require('pptxgenjs');

async function createPresentation() {{
    const pptx = new PptxGenJS();

    // Set properties
    pptx.layout = 'LAYOUT_{inp.aspect_ratio.replace(":", "_")}';
    pptx.title = '{inp.title}';
    pptx.author = '{inp.author}';

    // Add title slide
    const slide1 = pptx.addSlide();
    slide1.addText('{inp.title}', {{
        x: 1,
        y: 2,
        w: 8,
        h: 1.5,
        fontSize: 44,
        bold: true,
        align: 'center',
    }});

    // Add content slide
    const slide2 = pptx.addSlide();
    slide2.addText('Content Slide', {{
        x: 0.5,
        y: 0.5,
        w: 9,
        h: 1,
        fontSize: 32,
        bold: true,
    }});
    slide2.addText([
        {{ text: 'First bullet point', options: {{ bullet: true }} }},
        {{ text: 'Second bullet point', options: {{ bullet: true }} }},
        {{ text: 'Third bullet point', options: {{ bullet: true }} }},
    ], {{
        x: 0.5,
        y: 1.5,
        w: 9,
        h: 3,
        fontSize: 18,
    }});

    await pptx.writeFile('output.pptx');
}}

createPresentation();
"""
        else:  # basic
            code = '''"""Create basic PowerPoint using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(output_path: str, title: str, slides: list[dict]) -> None:
    """Create a basic PowerPoint presentation.

    Args:
        output_path: Path for output PPTX
        title: Presentation title
        slides: List of slide dicts with 'title' and 'content' keys
    """
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title

    # Content slides
    content_layout = prs.slide_layouts[1]
    for slide_data in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data.get('title', '')

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.text = slide_data.get('content', '')

    prs.save(output_path)
'''

        return json.dumps(
            {"code": code, "type": inp.creation_type, "aspect_ratio": inp.aspect_ratio}
        )
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def pptx_generate_slide_html(
    layout_type: str,
    aspect_ratio: str = "16:9",
    title: str = "",
    content: str = "",
    palette_name: str = "classic_blue",
) -> str:
    """Generate HTML for a slide to use with html2pptx. Returns HTML string."""
    try:
        inp = SlideHtmlInput(
            layout_type=layout_type,
            aspect_ratio=aspect_ratio,
            title=title,
            content=content,
            palette_name=palette_name,
        )

        palette = COLOR_PALETTES.get(inp.palette_name, COLOR_PALETTES["classic_blue"])
        html = _generate_slide_html(
            inp.layout_type, inp.aspect_ratio, inp.title, inp.content, palette
        )

        return json.dumps(
            {
                "html": html,
                "layout": inp.layout_type,
                "palette": palette["name"],
                "dimensions": _get_dimensions(inp.aspect_ratio),
            }
        )
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def pptx_generate_template_code(
    workflow_step: str,
) -> str:
    """Generate code for template-based PowerPoint workflow. Returns Python code string."""
    try:
        inp = TemplateInput(workflow_step=workflow_step)

        if inp.workflow_step == "full":
            code = TEMPLATE_WORKFLOW_CODE
        elif inp.workflow_step == "analyze":
            code = '''"""Analyze template presentation."""
import subprocess

def analyze_template(template_path: str, output_dir: str) -> dict:
    """Extract text and create thumbnails from template.

    Args:
        template_path: Path to template PPTX
        output_dir: Working directory

    Returns:
        Analysis results
    """
    # Extract text to markdown
    result = subprocess.run(
        ["python", "-m", "markitdown", template_path],
        capture_output=True,
        text=True,
    )

    # Create thumbnail grid
    subprocess.run(
        ["python", "scripts/thumbnail.py", template_path, f"{output_dir}/thumbnails", "--cols", "4"],
        check=True,
    )

    return {
        "text": result.stdout,
        "thumbnails": f"{output_dir}/thumbnails.jpg",
    }
'''
        elif inp.workflow_step == "rearrange":
            code = '''"""Rearrange slides from template."""
import subprocess

def rearrange_slides(template_path: str, output_path: str, slide_indices: list[int]) -> None:
    """Create new presentation with slides in specified order.

    Args:
        template_path: Path to template PPTX
        output_path: Path for output PPTX
        slide_indices: List of 0-based slide indices
    """
    order_str = ",".join(str(i) for i in slide_indices)
    subprocess.run(
        ["python", "scripts/rearrange.py", template_path, output_path, order_str],
        check=True,
    )
'''
        elif inp.workflow_step == "inventory":
            code = '''"""Extract text inventory from presentation."""
import subprocess
import json
from pathlib import Path

def extract_inventory(pptx_path: str, output_json: str) -> dict:
    """Extract all text shapes with positions and formatting.

    Args:
        pptx_path: Path to PPTX file
        output_json: Path for output JSON

    Returns:
        Inventory data structure
    """
    subprocess.run(
        ["python", "scripts/inventory.py", pptx_path, output_json],
        check=True,
    )

    return json.loads(Path(output_json).read_text())
'''
        else:  # replace
            code = '''"""Replace text in presentation."""
import subprocess

def replace_text(pptx_path: str, replacements_json: str, output_path: str) -> None:
    """Apply text replacements to presentation.

    Args:
        pptx_path: Path to input PPTX
        replacements_json: Path to JSON with replacement paragraphs
        output_path: Path for output PPTX
    """
    subprocess.run(
        ["python", "scripts/replace.py", pptx_path, replacements_json, output_path],
        check=True,
    )
'''

        return json.dumps({"code": code, "step": inp.workflow_step})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def pptx_generate_edit_code(
    edit_type: str,
) -> str:
    """Generate code for OOXML-based PowerPoint editing. Returns Python code string."""
    try:
        inp = EditInput(edit_type=edit_type)

        if inp.edit_type == "text":
            code = OOXML_EDIT_TEMPLATE
        elif inp.edit_type == "notes":
            code = '''"""Edit speaker notes via OOXML."""
import subprocess
from pathlib import Path
import xml.etree.ElementTree as ET

NAMESPACES = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

def get_notes(unpacked_dir: str, slide_num: int) -> str:
    """Get speaker notes for a slide.

    Args:
        unpacked_dir: Path to unpacked PPTX
        slide_num: Slide number (1-based)

    Returns:
        Notes text content
    """
    notes_path = Path(unpacked_dir) / f"ppt/notesSlides/notesSlide{slide_num}.xml"
    if not notes_path.exists():
        return ""

    tree = ET.parse(notes_path)
    root = tree.getroot()

    texts = []
    for t_elem in root.findall(".//a:t", NAMESPACES):
        if t_elem.text:
            texts.append(t_elem.text)

    return "\\n".join(texts)

def set_notes(unpacked_dir: str, slide_num: int, notes_text: str) -> None:
    """Set speaker notes for a slide.

    Args:
        unpacked_dir: Path to unpacked PPTX
        slide_num: Slide number (1-based)
        notes_text: New notes content
    """
    notes_path = Path(unpacked_dir) / f"ppt/notesSlides/notesSlide{slide_num}.xml"
    # Implementation requires creating/updating notes XML structure
    pass
'''
        else:  # shape or image
            code = '''"""Edit shapes/images via OOXML."""
import subprocess
from pathlib import Path
import xml.etree.ElementTree as ET

NAMESPACES = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

def list_shapes(unpacked_dir: str, slide_num: int) -> list[dict]:
    """List all shapes on a slide.

    Args:
        unpacked_dir: Path to unpacked PPTX
        slide_num: Slide number (1-based)

    Returns:
        List of shape info dicts
    """
    slide_path = Path(unpacked_dir) / f"ppt/slides/slide{slide_num}.xml"
    tree = ET.parse(slide_path)
    root = tree.getroot()

    shapes = []
    for sp in root.findall(".//p:sp", NAMESPACES):
        nvpr = sp.find(".//p:cNvPr", NAMESPACES)
        if nvpr is not None:
            shapes.append({
                "id": nvpr.get("id"),
                "name": nvpr.get("name"),
            })

    return shapes

def add_image(unpacked_dir: str, slide_num: int, image_path: str, x: int, y: int, width: int, height: int) -> None:
    """Add image to slide.

    Args:
        unpacked_dir: Path to unpacked PPTX
        slide_num: Slide number (1-based)
        image_path: Path to image file
        x, y: Position in EMUs
        width, height: Size in EMUs
    """
    # Copy image to media folder
    # Update relationships
    # Add pic element to slide XML
    pass
'''

        return json.dumps({"code": code, "type": inp.edit_type})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def pptx_generate_color_palette(
    mood: str = "",
    industry: str = "",
) -> str:
    """Generate color palette recommendations for presentations. Returns palette options."""
    try:
        inp = ColorPaletteInput(mood=mood, industry=industry)

        # Score palettes based on mood/industry keywords
        scored = []
        search_terms = f"{inp.mood} {inp.industry}".lower()

        for name, palette in COLOR_PALETTES.items():
            score = 0
            palette_mood = palette["mood"].lower()

            # Check for keyword matches
            for term in search_terms.split():
                if term in palette_mood or term in name:
                    score += 1

            scored.append((score, name, palette))

        # Sort by score descending
        scored.sort(key=lambda x: -x[0])

        # Return top 5 recommendations
        recommendations = []
        for score, name, palette in scored[:5]:
            recommendations.append(
                {
                    "name": palette["name"],
                    "id": name,
                    "colors": palette["colors"],
                    "mood": palette["mood"],
                    "match_score": score,
                }
            )

        return json.dumps(
            {
                "recommendations": recommendations,
                "search_terms": search_terms.strip(),
                "web_safe_fonts": WEB_SAFE_FONTS,
            }
        )
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def pptx_generate_layout(
    layout_type: str,
    aspect_ratio: str = "16:9",
) -> str:
    """Generate layout configuration for slides. Returns layout specs."""
    try:
        inp = LayoutInput(layout_type=layout_type, aspect_ratio=aspect_ratio)
        width, height = _get_dimensions(inp.aspect_ratio)

        layouts = {
            "title": {
                "name": "Title Slide",
                "regions": [
                    {
                        "name": "title",
                        "x": 50,
                        "y": height * 0.35,
                        "w": width - 100,
                        "h": 80,
                        "font_size": 48,
                    },
                    {
                        "name": "subtitle",
                        "x": 50,
                        "y": height * 0.55,
                        "w": width - 100,
                        "h": 40,
                        "font_size": 24,
                    },
                ],
            },
            "content": {
                "name": "Content Slide",
                "regions": [
                    {
                        "name": "header",
                        "x": 30,
                        "y": 20,
                        "w": width - 60,
                        "h": 50,
                        "font_size": 32,
                    },
                    {
                        "name": "body",
                        "x": 30,
                        "y": 80,
                        "w": width - 60,
                        "h": height - 110,
                        "font_size": 18,
                    },
                ],
            },
            "two_column": {
                "name": "Two Column",
                "regions": [
                    {
                        "name": "header",
                        "x": 30,
                        "y": 20,
                        "w": width - 60,
                        "h": 50,
                        "font_size": 32,
                    },
                    {
                        "name": "left",
                        "x": 30,
                        "y": 80,
                        "w": (width - 90) * 0.4,
                        "h": height - 110,
                        "font_size": 16,
                    },
                    {
                        "name": "right",
                        "x": 30 + (width - 90) * 0.45,
                        "y": 80,
                        "w": (width - 90) * 0.55,
                        "h": height - 110,
                        "font_size": 16,
                    },
                ],
            },
            "image_text": {
                "name": "Image + Text",
                "regions": [
                    {
                        "name": "text",
                        "x": 30,
                        "y": 20,
                        "w": (width - 60) * 0.5,
                        "h": height - 40,
                        "font_size": 18,
                    },
                    {
                        "name": "image",
                        "x": (width + 30) * 0.5,
                        "y": 20,
                        "w": (width - 60) * 0.5,
                        "h": height - 40,
                        "type": "placeholder",
                    },
                ],
            },
            "quote": {
                "name": "Quote Slide",
                "regions": [
                    {
                        "name": "quote",
                        "x": 60,
                        "y": height * 0.3,
                        "w": width - 120,
                        "h": 100,
                        "font_size": 28,
                        "style": "italic",
                    },
                    {
                        "name": "attribution",
                        "x": 60,
                        "y": height * 0.65,
                        "w": width - 120,
                        "h": 30,
                        "font_size": 18,
                    },
                ],
            },
            "section": {
                "name": "Section Divider",
                "regions": [
                    {
                        "name": "title",
                        "x": 50,
                        "y": height * 0.4,
                        "w": width - 100,
                        "h": 80,
                        "font_size": 56,
                        "align": "center",
                    },
                ],
            },
        }

        layout = layouts.get(inp.layout_type, layouts["content"])

        return json.dumps(
            {
                "layout": layout,
                "dimensions": {"width": width, "height": height, "unit": "pt"},
                "aspect_ratio": inp.aspect_ratio,
            }
        )
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def pptx_generate_thumbnail_code(
    include_pdf_conversion: bool = False,
) -> str:
    """Generate code for creating thumbnails and visual validation. Returns Python code."""
    try:
        inp = ThumbnailInput(include_pdf_conversion=include_pdf_conversion)
        code = THUMBNAIL_CODE

        if not inp.include_pdf_conversion:
            # Remove PDF conversion functions
            lines = code.split("\n")
            kept = []
            skip_until_def = False

            for line in lines:
                if "def convert_to_pdf" in line or "def pdf_to_images" in line:
                    skip_until_def = True
                    continue
                if skip_until_def and line.startswith("def "):
                    skip_until_def = False
                if not skip_until_def:
                    kept.append(line)

            code = "\n".join(kept)

        return json.dumps({"code": code, "include_pdf": inp.include_pdf_conversion})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def pptx_detect_antipatterns(
    code: str,
) -> str:
    """Detect anti-patterns in PPTX/HTML code. Returns findings as JSON."""
    try:
        inp = AntipatternInput(code=code)
        findings = _detect_antipatterns(inp.code)

        return json.dumps(
            {
                "findings": findings,
                "count": len(findings),
                "clean": len(findings) == 0,
            }
        )
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
async def pptx_generate_scaffold(
    project_name: str,
    workflow: str = "html2pptx",
) -> str:
    """Generate complete PPTX project scaffold. Returns file structure as JSON."""
    try:
        inp = ScaffoldInput(project_name=project_name, workflow=workflow)

        files = {
            f"{inp.project_name}/README.md": f"""# {inp.project_name}

PowerPoint presentation project using {inp.workflow} workflow.

## Setup

```bash
npm install pptxgenjs playwright
pip install markitdown defusedxml
```

## Usage

See `create_presentation.js` for the main workflow.
""",
        }

        if inp.workflow == "html2pptx":
            files[f"{inp.project_name}/create_presentation.js"] = (
                HTML2PPTX_TEMPLATE.format(
                    aspect_ratio="16_9",
                    title=inp.project_name,
                    author="",
                )
            )
            files[f"{inp.project_name}/slides/slide1.html"] = _generate_slide_html(
                "title",
                "16:9",
                "Presentation Title",
                "Subtitle goes here",
                COLOR_PALETTES["classic_blue"],
            )
        elif inp.workflow == "template":
            files[f"{inp.project_name}/workflow.py"] = TEMPLATE_WORKFLOW_CODE
        else:  # ooxml
            files[f"{inp.project_name}/edit.py"] = OOXML_EDIT_TEMPLATE

        files[f"{inp.project_name}/validate.py"] = THUMBNAIL_CODE

        return json.dumps(
            {
                "project": inp.project_name,
                "workflow": inp.workflow,
                "files": files,
            }
        )
    except Exception as e:
        return json.dumps({"error": str(e)})


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    mcp.run(transport="stdio")
